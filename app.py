# app.py

import streamlit as st
import google.generativeai as genai
import time
import re
import os
import io

# Import tools for file processing
import fitz  # PyMuPDF
import docx
import pptx
import markdown2
import pdfkit
import graphviz
from PIL import Image

# --- Page Configuration ---
st.set_page_config(
    page_title="AI Study Pack Generator",
    page_icon="🧠",
    layout="wide"
)

# --- App Header ---
st.title("🧠 AI Study Pack Generator")

# --- API Key Logic ---
api_key = ""
try:
    # Attempt to load the key from Streamlit's secrets
    api_key = st.secrets["GEMINI_API_KEY"]
    st.success("API key successfully loaded!", icon="✅")
except KeyError:
    # If not found in secrets, show a red cross and an input box
    st.error("API key not found. Please add your key below to proceed.", icon="❌")
    api_key = st.text_input("Enter your Gemini API Key here:", type="password")

# --- Main App Logic (only runs if an API key is provided) ---
if api_key:
    genai.configure(api_key=api_key)

    st.write("Upload your class notes (.pdf, .docx, .pptx, images) and get a complete study package: comprehensive notes and a concept mind map!")
    
    # --- Advanced Options (moved from sidebar to an expander) ---
    with st.expander("Advanced Options"):
        system_message = st.text_area(
            "System Message (AI's Role)",
            "You are an expert academic assistant. Your goal is to transform notes into a comprehensive, well-structured educational document. Elaborate on topics, add examples, and clarify complex concepts in Markdown format."
        )

    # --- File Uploader ---
    uploaded_files = st.file_uploader(
        "Upload your notes here",
        type=['pdf', 'docx', 'pptx', 'png', 'jpg', 'jpeg'],
        accept_multiple_files=True
    )

    # --- Generate Button ---
    if st.button("Generate Study Pack", type="primary", use_container_width=True):
        if uploaded_files:
            # (The rest of the processing logic remains the same)
            # --- 1. Content Extraction ---
            with st.status("Step 1/4: Extracting content from files...", expanded=True) as status:
                final_extracted_content = ""
                images_for_processing = []

                for uploaded_file in uploaded_files:
                    st.write(f"Processing: {uploaded_file.name}")
                    if uploaded_file.type == "application/pdf":
                        pdf_bytes = uploaded_file.getvalue()
                        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
                        for page in pdf_document:
                            pix = page.get_pixmap(dpi=150)
                            img = Image.open(io.BytesIO(pix.tobytes("png")))
                            images_for_processing.append(img)
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                        doc = docx.Document(uploaded_file)
                        final_extracted_content += "\n".join([p.text for p in doc.paragraphs])
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                        prs = pptx.Presentation(uploaded_file)
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    final_extracted_content += shape.text + "\n"
                    else: # Handle images
                        images_for_processing.append(Image.open(uploaded_file))
                
                if images_for_processing:
                    st.write("Images detected, preparing for vision analysis...")
                    model_input = [final_extracted_content] + images_for_processing
                    vision_model = genai.GenerativeModel('gemini-2.5-pro')
                    vision_response = vision_model.generate_content(model_input)
                    final_extracted_content = vision_response.text
                
                st.write("✅ Content extraction complete.")
                status.update(label="Content extracted successfully!", state="complete")

            # --- 2. Generate Comprehensive Notes ---
            with st.status("Step 2/4: Generating comprehensive notes with AI...", expanded=True) as status:
                try:
                    text_model = genai.GenerativeModel(model_name='gemini-2.5-pro', system_instruction=system_message)
                    notes_prompt = f"Please take the following extracted class notes and expand them into a detailed, comprehensive document. For each key point, provide detailed explanations, include concrete examples, and organize the information logically.\n\n---\n{final_extracted_content}\n---"
                    final_response = text_model.generate_content(notes_prompt, request_options={'timeout': 600})
                    generated_notes = final_response.text
                    st.write("✅ Comprehensive notes generated.")
                    status.update(label="Notes generated!", state="complete")
                except Exception as e:
                    st.error(f"An error occurred while generating notes: {e}")
                    st.stop()
            
            # --- 3. Generate Mind Map ---
            with st.status("Step 3/4: Creating concept mind map...", expanded=True) as status:
                st.write("🗺️ Generating mind map...")
                time.sleep(20)
                mindmap_prompt = f"Analyze the following text and generate a structural mind map in Graphviz DOT language. The central node should be the main topic. Create clear, hierarchical relationships. Keep node labels concise (1-3 words). Do not include any explanation, only the DOT code itself.\nText:\n---\n{generated_notes}"
                mindmap_response = text_model.generate_content(mindmap_prompt)
                dot_code = re.sub(r'```dot\s*|```', '', mindmap_response.text).strip()
                
                st.write("✅ Mind map structure created.")
                status.update(label="Mind map created!", state="complete")

            # --- 4. Prepare and Display Downloads ---
            with st.status("Step 4/4: Preparing your download package...", expanded=True) as status:
                html_text = markdown2.markdown(generated_notes, extras=["tables", "fenced-code-blocks", "code-friendly"])
                html_with_style = f'<html><head><meta charset="utf-8"><style>body{{font-family: Arial, sans-serif;}} pre, code {{background-color: #f4f4f4; padding: 1em; border-radius: 5px;}}</style></head><body>{html_text}</body></html>'
                pdfkit.from_string(html_with_style, "Generated_Notes.pdf", options={"enable-local-file-access": ""})
                
                src = graphviz.Source(dot_code)
                src.render("Concept_Mind_Map", format='png', view=False, cleanup=True)
                
                st.write("✅ Package ready for download.")
                status.update(label="Downloads are ready!", state="complete")
            
            # --- Display Results and Download Buttons ---
            st.header("Your Study Pack is Ready!", divider="rainbow")
            col1, col2 = st.columns(2)
            with col1:
                st.subheader("📝 Generated Notes (PDF)")
                with open("Generated_Notes.pdf", "rb") as file:
                    st.download_button(label="Download Notes PDF", data=file, file_name="Generated_Notes.pdf", mime="application/pdf", use_container_width=True)
            with col2:
                st.subheader("🗺️ Concept Mind Map")
                st.image("Concept_Mind_Map.png")
                with open("Concept_Mind_Map.png", "rb") as file:
                    st.download_button(label="Download Mind Map PNG", data=file, file_name="Concept_Mind_Map.png", mime="image/png", use_container_width=True)
        else:
            st.warning("Please upload your notes to get started.")
else:
    # This message shows if the API key is not in secrets and the user hasn't entered one yet.
    st.info("Please provide an API key to use the app.")